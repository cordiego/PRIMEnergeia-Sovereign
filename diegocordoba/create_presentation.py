import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define colors matching the LaTeX theme
primeblue = RGBColor(0x0A, 0x16, 0x28)
primegold = RGBColor(0xD4, 0xA0, 0x17)
primegray = RGBColor(0xF4, 0xF6, 0xF9)
primegreen = RGBColor(0x2D, 0x6A, 0x4F)

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "PRIMEnergeia S.A.S."
title.text_frame.paragraphs[0].font.color.rgb = primeblue
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "INTELIGENCIA SOBERANA PARA INFRAESTRUCTURA ENERGÉTICA CRÍTICA\n\nControl Estocástico HJB · Perovskita de Alta Eficiencia · Arbitraje PML · SCADA/Modbus/DNP3"
subtitle.text_frame.paragraphs[0].font.color.rgb = primegold

# Slide 2: La Tesis (Value Proposition)
bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
shapes = slide2.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Propuesta de Valor: La Tesis"
title_shape.text_frame.paragraphs[0].font.color.rgb = primeblue

tf = body_shape.text_frame
tf.text = "Los activos renovables en México pierden valor silenciosamente debido a:"
p = tf.add_paragraph()
p.text = "Congestión nodal"
p.level = 1
p = tf.add_paragraph()
p.text = "Entropía térmica en transformadores"
p.level = 1
p = tf.add_paragraph()
p.text = "PML volátil"
p.level = 1

p2 = tf.add_paragraph()
p2.text = "Estas generan pérdidas sistemáticas que los modelos determinísticos convencionales no pueden capturar ni mitigar."
p2.level = 0
p3 = tf.add_paragraph()
p3.text = "PRIMEnergeia resuelve esto en tiempo real."
p3.level = 0
p3.font.color.rgb = primegreen
p3.font.bold = True

# Slide 3: El Problema
slide3 = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide3.shapes.title
body_shape = slide3.shapes.placeholders[1]

title_shape.text = "El Problema"
title_shape.text_frame.paragraphs[0].font.color.rgb = primeblue

tf = body_shape.text_frame
tf.text = "Frecuencia nodal inestable fuera de ±0.2 Hz genera penalizaciones CENACE."
p = tf.add_paragraph()
p.text = "Entropía térmica acelerada en transformadores: ciclo de vida reducido 18--34%."
p = tf.add_paragraph()
p.text = "Arbitraje PML desaprovechado en nodos 400 kV: pérdida anual estimada $80K--$160K USD por nodo."
p = tf.add_paragraph()
p.text = "Módulos solares convencionales: 0% reciclaje de fotones, micro-grietas no contenidas."

# Set color to red for problem bullets
for paragraph in tf.paragraphs:
    paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

# Slide 4: La Solución
slide4 = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide4.shapes.title
body_shape = slide4.shapes.placeholders[1]

title_shape.text = "La Solución"
title_shape.text_frame.paragraphs[0].font.color.rgb = primeblue

tf = body_shape.text_frame
tf.text = "HJB Solver — resolución en tiempo real de Hamilton-Jacobi-Bellman."
p = tf.add_paragraph()
p.text = "DRL Auto-Healing — Actor-Critic para recuperación post-disturbancia."
p = tf.add_paragraph()
p.text = "Granas Module — geometría propietaria 21×34, 89% photon recycling."
p = tf.add_paragraph()
p.text = "SIBO Optimizer — Bayesiano para optimización continua de parámetros."

# Set color to primegreen for solution bullets
for paragraph in tf.paragraphs:
    paragraph.font.color.rgb = primegreen

# Save presentation
prs.save('/Users/diegocordoba/PRIMEnergeia_Pitch.pptx')
