import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# Set 16:9 aspect ratio for modern look
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
primeblue = RGBColor(0x0A, 0x16, 0x28)
primegold = RGBColor(0xD4, 0xA0, 0x17)
white = RGBColor(0xFF, 0xFF, 0xFF)
light_gray = RGBColor(0xE0, 0xE0, 0xE0)
green_accent = RGBColor(0x2A, 0x9D, 0x8F)

blank_slide_layout = prs.slide_layouts[6] # Blank layout for full control

def add_dark_background(slide):
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = primeblue
    background.line.color.rgb = primeblue

# ----------------- SLIDE 1: Title -----------------
slide1 = prs.slides.add_slide(blank_slide_layout)
add_dark_background(slide1)

txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.5), prs.slide_width - Inches(2), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "PRIMEnergeia"
p.font.size = Pt(90)
p.font.color.rgb = white
p.font.bold = True

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "INTELIGENCIA SOBERANA PARA INFRAESTRUCTURA ENERGÉTICA CRÍTICA"
p2.font.size = Pt(22)
p2.font.color.rgb = primegold
p2.font.bold = True
p2.space_before = Pt(20)

# ----------------- SLIDE 2: Thesis -----------------
slide2 = prs.slides.add_slide(blank_slide_layout)
add_dark_background(slide2)

txBox = slide2.shapes.add_textbox(Inches(1.5), Inches(2.5), prs.slide_width - Inches(3), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "Los activos renovables pierden valor silenciosamente."
p.font.size = Pt(40)
p.font.color.rgb = light_gray

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "PRIMEnergeia resuelve esto en tiempo real."
p2.font.size = Pt(48)
p2.font.color.rgb = primegold
p2.font.bold = True
p2.space_before = Pt(30)

# ----------------- SLIDE 3: Problem -----------------
slide3 = prs.slides.add_slide(blank_slide_layout)
add_dark_background(slide3)

title_box = slide3.shapes.add_textbox(Inches(1.5), Inches(1), Inches(10), Inches(1))
title_tf = title_box.text_frame
p = title_tf.paragraphs[0]
p.text = "EL PROBLEMA"
p.font.size = Pt(32)
p.font.color.rgb = primegold
p.font.bold = True

content_box = slide3.shapes.add_textbox(Inches(1.5), Inches(2.5), prs.slide_width - Inches(3), Inches(4))
tf = content_box.text_frame
tf.word_wrap = True

problems = [
    "Frecuencia nodal inestable (fuera de ±0.2 Hz)",
    "Entropía térmica acelerada en transformadores",
    "Arbitraje PML desaprovechado (pérdidas $80K-$160K)",
    "Módulos solares convencionales ineficientes"
]

for prob in problems:
    p = tf.add_paragraph()
    p.text = "•  " + prob
    p.font.size = Pt(32)
    p.font.color.rgb = light_gray
    p.space_after = Pt(24)

# ----------------- SLIDE 4: Solution -----------------
slide4 = prs.slides.add_slide(blank_slide_layout)
add_dark_background(slide4)

title_box = slide4.shapes.add_textbox(Inches(1.5), Inches(1), Inches(10), Inches(1))
title_tf = title_box.text_frame
p = title_tf.paragraphs[0]
p.text = "LA SOLUCIÓN"
p.font.size = Pt(32)
p.font.color.rgb = green_accent
p.font.bold = True

content_box = slide4.shapes.add_textbox(Inches(1.5), Inches(2.5), prs.slide_width - Inches(3), Inches(4))
tf = content_box.text_frame
tf.word_wrap = True

solutions = [
    "HJB Solver: Control estocástico en tiempo real",
    "DRL Auto-Healing: Actor-Critic post-disturbancia",
    "Granas Module: Geometría 21×34, 89% photon recycling",
    "SIBO Optimizer: Calibración Bayesiana continua"
]

for sol in solutions:
    p = tf.add_paragraph()
    p.text = "•  " + sol
    p.font.size = Pt(32)
    p.font.color.rgb = white
    p.space_after = Pt(24)

prs.save('/Users/diegocordoba/PRIMEnergeia_Pitch_Minimal.pptx')
